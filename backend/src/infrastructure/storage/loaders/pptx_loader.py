from pathlib import Path
from typing import List
import pptx
from backend.src.infrastructure.storage.loaders.base_loader import BaseDocumentLoader, ExtractedChunk


class PPTXDocumentLoader(BaseDocumentLoader):
    """Loader strategy for extracting text from PPTX files using python-pptx."""

    def load(self, file_path: Path) -> List[ExtractedChunk]:
        presentation = pptx.Presentation(str(file_path))
        chunks: List[ExtractedChunk] = []

        for slide_idx, slide in enumerate(presentation.slides, start=1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())

            slide_content = "\n".join(slide_texts)
            if slide_content.strip():
                chunks.append(
                    ExtractedChunk(
                        content=slide_content,
                        page_number=slide_idx,
                        metadata={"slide_number": slide_idx, "total_slides": len(presentation.slides)},
                    )
                )

        return chunks
